#!/usr/bin/env python3
"""Generate SIX AI presentation as PPTX matching the HTML version exactly."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Slide dimensions (16:9) ────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)

# ── Colors ──────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT    = RGBColor(0xF8, 0xFA, 0xFC)
SURFACE     = RGBColor(0xF1, 0xF5, 0xF9)
BORDER      = RGBColor(0xCB, 0xD5, 0xE1)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
TEXT_DARK   = RGBColor(0x0F, 0x17, 0x2A)
TEXT_SUB    = RGBColor(0x47, 0x55, 0x69)
SIX_RED     = RGBColor(0xC8, 0x10, 0x2E)
SIX_RED_DK  = RGBColor(0x9B, 0x0C, 0x23)
SIX_BLUE    = RGBColor(0x2E, 0x6F, 0xD6)
SIX_BLUE_BR = RGBColor(0x4F, 0x8B, 0xF0)
NAVY        = RGBColor(0x14, 0x23, 0x3F)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)

RED_TINT    = RGBColor(0xFE, 0xF7, 0xF8)
BLUE_TINT   = RGBColor(0xEB, 0xF1, 0xFB)
GREEN_TINT  = RGBColor(0xE8, 0xF8, 0xEE)
AMBER_TINT  = RGBColor(0xFD, 0xF3, 0xE4)
PURPLE_TINT = RGBColor(0xF1, 0xEB, 0xFD)

DIR = os.path.dirname(os.path.abspath(__file__))

# ── Helpers ─────────────────────────────────────────────
def tb(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box.text_frame

def run(para, text, size=14, bold=False, color=TEXT_DARK, italic=False, name='Segoe UI'):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    r.font.name = name
    return r

def para(tf, text, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
         space_before=0, space_after=0, italic=False):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run(p, text, size, bold, color, italic)
    return p

def rect(slide, left, top, width, height, fill=WHITE, line_color=BORDER, line_w=Pt(1), radius=0.015):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line_color
    s.line.width = line_w
    if len(s.adjustments) > 0:
        s.adjustments[0] = radius
    return s

def pill(slide, left, top, text, bg, fg, w=None, h=Inches(0.38), size=12):
    w = w or Inches(max(1.5, len(text) * 0.135))
    s = rect(slide, left, top, w, h, bg, fg, Pt(0.75), 0.12)
    tf = s.text_frame
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, text, size, False, fg)
    return s

def topbar(slide, label_text):
    """SIX AI text-logo top-left, label center, SwissHacks logo top-right."""
    # Top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.55), SW, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()

    # SIX logo image + AI text
    six_logo = os.path.join(DIR, 'six_logo_small.png')
    if os.path.exists(six_logo):
        slide.shapes.add_picture(six_logo, Inches(0.45), Inches(0.15), Inches(0.55), Inches(0.28))
    tf = tb(slide, Inches(1.0), Inches(0.18), Inches(0.5), Inches(0.3))
    p = tf.paragraphs[0]
    run(p, 'AI', 14, False, TEXT_SUB)

    # Center label
    tf2 = tb(slide, Inches(4), Inches(0.2), Inches(5.333), Inches(0.3))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run(p2, label_text, 10, False, MUTED)

    # SwissHacks logo
    logo = os.path.join(DIR, 'pptx_slide2_Google_Shape_64_p14.png')
    if os.path.exists(logo):
        slide.shapes.add_picture(logo, Inches(11.5), Inches(0.12), Inches(1.5))


# ══════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]

# Margin constants
LM = Inches(0.8)   # left margin
CONTENT_W = Inches(11.733)  # 13.333 - 2*0.8
CENTER = PP_ALIGN.CENTER


# ══════════════════════════════════════════════════════════
# SLIDE 0 — TITLE
# ══════════════════════════════════════════════════════════
s0 = prs.slides.add_slide(blank)
s0.background.fill.solid()
s0.background.fill.fore_color.rgb = BG_LIGHT

# SwissHacks logo top-right
logo = os.path.join(DIR, 'pptx_slide2_Google_Shape_64_p14.png')
if os.path.exists(logo):
    s0.shapes.add_picture(logo, Inches(11.5), Inches(0.3), Inches(1.5))

# SIX AI large title — logo image + AI text
six_logo_lg = os.path.join(DIR, 'six_logo_large.png')
if os.path.exists(six_logo_lg):
    logo_w = Inches(2.4)
    logo_h = Inches(0.67)
    logo_x = Inches((13.333 - 2.4 - 0.2 - 1.2) / 2)
    s0.shapes.add_picture(six_logo_lg, logo_x, Inches(2.0), logo_w, logo_h)
    tf = tb(s0, Inches(logo_x.inches + 2.4 + 0.15), Inches(1.75), Inches(1.5), Inches(1.2))
    p = tf.paragraphs[0]
    run(p, 'AI', 54, False, TEXT_DARK)

# Quote
tf2 = tb(s0, Inches(2), Inches(2.8), Inches(9.333), Inches(1.2))
p2 = tf2.paragraphs[0]
p2.alignment = CENTER
run(p2, '"Intellectuals solve problems, geniuses prevent them."', 22, False, TEXT_SUB, italic=True)
p3 = tf2.add_paragraph()
p3.alignment = CENTER
p3.space_before = Pt(8)
run(p3, '— Albert Einstein, probably', 14, False, MUTED)

# Team names with dividers
names = ['Boris Goranov', 'Hristo Stefanov', 'Nikola Staykov']
total_w = Inches(8)
start_x = Inches((13.333 - 8) / 2)
name_y = Inches(4.5)
name_w = Inches(2.4)
div_gap = Inches(0.2)

for i, name in enumerate(names):
    x = start_x + i * (name_w + div_gap + Pt(1) + div_gap)
    tf_n = tb(s0, x, name_y, name_w, Inches(0.4))
    para(tf_n, name, 16, True, TEXT_SUB, CENTER)
    if i < 2:
        dx = x + name_w + div_gap
        divider = s0.shapes.add_shape(MSO_SHAPE.RECTANGLE, dx, name_y + Inches(0.05), Pt(1), Inches(0.3))
        divider.fill.solid()
        divider.fill.fore_color.rgb = BORDER
        divider.line.fill.background()

# SwissHacks 2026
tf3 = tb(s0, Inches(3), Inches(5.8), Inches(7.333), Inches(0.4))
para(tf3, 'SWISSHACKS 2026', 12, True, MUTED, CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 1 — PROBLEM
# ══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
s1.background.fill.solid()
s1.background.fill.fore_color.rgb = WHITE
topbar(s1, '02 / 06 — PROBLEM')

# Eyebrow
tf = tb(s1, LM, Inches(0.9), CONTENT_W, Inches(0.3))
para(tf, 'THE PROBLEM', 12, True, SIX_RED, CENTER)

# Title
tf2 = tb(s1, Inches(1.2), Inches(1.3), Inches(10.933), Inches(1.2))
p = tf2.paragraphs[0]
p.alignment = CENTER
run(p, 'Private banking is personal — ', 36, True, TEXT_DARK)
run(p, 'but only if you have CHF 50M+', 36, True, SIX_RED, italic=True)

# Subtitle
tf3 = tb(s1, Inches(2), Inches(2.6), Inches(9.333), Inches(0.7))
para(tf3, 'Ultra-high-net-worth clients get a dedicated RM who knows them deeply.\nEveryone else gets a template.', 16, False, TEXT_SUB, CENTER)

# Pain cards — icon + title side by side, centered
cards = [
    ('🏦', 'The UHNW privilege gap'),
    ('⏱️', "RMs can't scale intimacy"),
    ('📉', 'Generic service loses clients'),
]

card_w = Inches(3.6)
card_h = Inches(1.0)
card_gap = Inches(0.4)
total_cards_w = 3 * card_w.inches + 2 * card_gap.inches
card_start_x = Inches((13.333 - total_cards_w) / 2)
card_y = Inches(4.0)

for i, (icon, title) in enumerate(cards):
    x = Inches(card_start_x.inches + i * (card_w.inches + card_gap.inches))
    box = rect(s1, x, card_y, card_w, card_h)

    # Icon
    tf_i = tb(s1, x + Inches(0.35), card_y + Inches(0.15), Inches(0.6), Inches(0.7))
    para(tf_i, icon, 28, False, TEXT_DARK, CENTER)

    # Title
    tf_t = tb(s1, x + Inches(1.0), card_y + Inches(0.2), card_w - Inches(1.3), Inches(0.6))
    tf_t.paragraphs[0].alignment = PP_ALIGN.LEFT
    run(tf_t.paragraphs[0], title, 16, True, TEXT_DARK)


# ══════════════════════════════════════════════════════════
# SLIDE 2 — SOLUTION
# ══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
s2.background.fill.solid()
s2.background.fill.fore_color.rgb = WHITE
topbar(s2, '03 / 06 — SOLUTION')

# Eyebrow
tf = tb(s2, LM, Inches(0.85), CONTENT_W, Inches(0.3))
para(tf, 'OUR SOLUTION', 12, True, SIX_BLUE, CENTER)

# Title
tf2 = tb(s2, Inches(1.2), Inches(1.2), Inches(10.933), Inches(0.8))
p = tf2.paragraphs[0]
p.alignment = CENTER
run(p, 'Give ', 36, True, TEXT_DARK)
run(p, 'every client', 36, True, SIX_BLUE, italic=True)
run(p, ' the UHNW experience — at scale', 36, True, TEXT_DARK)

# Flow pipeline
flow = [
    ('CRM History', SIX_BLUE, BLUE_TINT),
    ('Client DNA', SIX_RED, RED_TINT),
    ('Live Intelligence', PURPLE, PURPLE_TINT),
    ('Personalised Advisory', GREEN, GREEN_TINT),
]
flow_y = Inches(2.2)
fpill_w = Inches(2.2)
arrow_w = Inches(0.5)
total_flow_w = 4 * fpill_w.inches + 3 * arrow_w.inches
flow_start = Inches((13.333 - total_flow_w) / 2)

for i, (label, fg, bg) in enumerate(flow):
    x = Inches(flow_start.inches + i * (fpill_w.inches + arrow_w.inches))
    pill(s2, x, flow_y, label, bg, fg, w=fpill_w, h=Inches(0.4), size=13)
    if i < 3:
        tf_a = tb(s2, x + fpill_w, flow_y, arrow_w, Inches(0.4))
        para(tf_a, '→', 16, False, MUTED, CENTER)

# Feature grid 2x2
features = [
    ('🧬', 'AI Client DNA', 'Synthesises CRM history into a living profile — values, risk appetite, life events — in seconds, not decades.', RED_TINT, SIX_RED),
    ('⚡', 'Conflict & Mandate Alerts', 'Continuously checks every position against personal values and CIO ratings — flags misalignments before they cost trust.', AMBER_TINT, AMBER),
    ('📰', 'News Scored Per Client', "Market news ranked by relevance to each client's portfolio and beliefs, not generic asset-class feeds.", BLUE_TINT, SIX_BLUE),
    ('✉️', 'Advisory in One Click', 'Tone-matched, context-aware client message in under 5 seconds — outreach that used to take an hour.', GREEN_TINT, GREEN),
]

feat_w = Inches(5.8)
feat_h = Inches(1.25)
feat_gap_x = Inches(0.4)
feat_gap_y = Inches(0.3)
feat_total_w = 2 * feat_w.inches + feat_gap_x.inches
feat_start_x = Inches((13.333 - feat_total_w) / 2)
feat_start_y = Inches(2.95)

for i, (icon, title, desc, tint, color) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(feat_start_x.inches + col * (feat_w.inches + feat_gap_x.inches))
    y = Inches(feat_start_y.inches + row * (feat_h.inches + feat_gap_y.inches))

    box = rect(s2, x, y, feat_w, feat_h)

    # Icon circle
    icon_box = rect(s2, x + Inches(0.2), y + Inches(0.2), Inches(0.55), Inches(0.55), tint, tint, Pt(0))
    itf = icon_box.text_frame
    itf.margin_top = Pt(2)
    p_ic = itf.paragraphs[0]
    p_ic.alignment = CENTER
    run(p_ic, icon, 18, False, TEXT_DARK)

    # Title
    tf_t = tb(s2, x + Inches(0.9), y + Inches(0.15), feat_w - Inches(1.1), Inches(0.35))
    para(tf_t, title, 15, True, TEXT_DARK)

    # Desc
    tf_d = tb(s2, x + Inches(0.9), y + Inches(0.5), feat_w - Inches(1.1), Inches(0.7))
    para(tf_d, desc, 11, False, TEXT_SUB)


# ══════════════════════════════════════════════════════════
# SLIDE 3 — ARCHITECTURE
# ══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
s3.background.fill.solid()
s3.background.fill.fore_color.rgb = WHITE
topbar(s3, '04 / 06 — ARCHITECTURE')

# Eyebrow
tf = tb(s3, LM, Inches(0.85), CONTENT_W, Inches(0.3))
para(tf, 'ARCHITECTURE', 12, True, SIX_RED, CENTER)

# Title
tf2 = tb(s3, Inches(2), Inches(1.2), Inches(9.333), Inches(0.7))
p = tf2.paragraphs[0]
p.alignment = CENTER
run(p, 'One ', 36, True, TEXT_DARK)
run(p, 'intelligent', 36, True, SIX_RED, italic=True)
run(p, ' platform', 36, True, TEXT_DARK)

# Subtitle
tf3 = tb(s3, Inches(1.5), Inches(1.9), Inches(10.333), Inches(0.6))
para(tf3, 'Real-time data from SIX Financial & Event Registry, processed by specialised LLM assistants, served through an RM-first interface.', 13, False, TEXT_SUB, CENTER)

# ── Data Sources row ──
ds_y = Inches(2.7)
ds_box = rect(s3, Inches(0.8), ds_y, Inches(11.733), Inches(0.95), SURFACE, BORDER)
tf_ds = tb(s3, Inches(1.1), ds_y + Inches(0.08), Inches(3), Inches(0.3))
p_ds = tf_ds.paragraphs[0]
run(p_ds, '●', 10, True, SIX_BLUE)
run(p_ds, '  DATA SOURCES', 11, True, SIX_BLUE)

ds_pills = ['SIX Financial MCP', 'Event Registry News', 'CRM Notes', 'Portfolio Data']
px = Inches(1.1)
for pt in ds_pills:
    pw = Inches(max(1.5, len(pt) * 0.14))
    pill(s3, px, ds_y + Inches(0.45), pt, BLUE_TINT, SIX_BLUE, w=pw, size=12)
    px += pw + Inches(0.2)

# Arrow
tf_arr = tb(s3, Inches(5), Inches(3.7), Inches(3.333), Inches(0.35))
p_arr = tf_arr.paragraphs[0]
p_arr.alignment = CENTER
run(p_arr, '↓  Live data feeds', 11, False, MUTED)

# ── Assistant/Agent row ──
agent_y = Inches(4.1)
agent_w = Inches(3.78)
agent_gap = Inches(0.2)
agent_start = Inches(0.8)

agents = [
    ('Conflict Assistant', 'Monitors the portfolio'),
    ('Message Assistant', 'Drafts the advisory notes'),
    ('Chat Agent', 'Natural language tool use across the system'),
]

for i, (name, desc) in enumerate(agents):
    x = Inches(agent_start.inches + i * (agent_w.inches + agent_gap.inches))
    box = rect(s3, x, agent_y, agent_w, Inches(0.9), RED_TINT, SIX_RED, Pt(1))

    tf_n = tb(s3, x + Inches(0.2), agent_y + Inches(0.08), agent_w - Inches(0.4), Inches(0.3))
    p_n = tf_n.paragraphs[0]
    run(p_n, '●', 10, True, SIX_RED)
    run(p_n, '  ' + name.upper(), 11, True, SIX_RED)

    pw = Inches(min(agent_w.inches - 0.4, max(1.8, len(desc) * 0.13)))
    pill(s3, x + Inches(0.2), agent_y + Inches(0.48), desc, RED_TINT, SIX_RED, w=pw, size=12)

# Arrow
tf_arr2 = tb(s3, Inches(5), Inches(5.05), Inches(3.333), Inches(0.35))
p_arr2 = tf_arr2.paragraphs[0]
p_arr2.alignment = CENTER
run(p_arr2, '↓  REST API', 11, False, MUTED)

# ── Bottom row ──
bot_y = Inches(5.5)
bot_h = Inches(1.25)

# Trust & Compliance
tc_w = Inches(5.6)
tc_box = rect(s3, Inches(0.8), bot_y, tc_w, bot_h, SURFACE, BORDER)
tf_tc = tb(s3, Inches(1.1), bot_y + Inches(0.08), Inches(3), Inches(0.3))
p_tc = tf_tc.paragraphs[0]
run(p_tc, '●', 10, True, GREEN)
run(p_tc, '  TRUST & COMPLIANCE', 11, True, GREEN)

tc_pills = ['Audit Trail', 'Explainability', 'Trace Logging']
px = Inches(1.1)
for pt in tc_pills:
    pw = Inches(max(1.2, len(pt) * 0.13))
    pill(s3, px, bot_y + Inches(0.45), pt, GREEN_TINT, GREEN, w=pw, size=11)
    px += pw + Inches(0.15)

pill(s3, Inches(1.1), bot_y + Inches(0.88), 'Spider Graph Scoring', GREEN_TINT, GREEN, w=Inches(2.2), size=11)

# RM Dashboard
rm_x = Inches(6.6)
rm_w = Inches(5.933)
rm_box = rect(s3, rm_x, bot_y, rm_w, bot_h, SURFACE, BORDER)
tf_rm = tb(s3, rm_x + Inches(0.3), bot_y + Inches(0.08), Inches(3), Inches(0.3))
p_rm = tf_rm.paragraphs[0]
run(p_rm, '●', 10, True, SIX_BLUE)
run(p_rm, '  RM DASHBOARD', 11, True, SIX_BLUE)

rm_pills = ['Home & Action Items', 'Global News View']
px = rm_x + Inches(0.3)
for pt in rm_pills:
    pw = Inches(max(1.5, len(pt) * 0.14))
    pill(s3, px, bot_y + Inches(0.45), pt, BLUE_TINT, SIX_BLUE, w=pw, size=11)
    px += pw + Inches(0.15)

pill(s3, rm_x + Inches(0.3), bot_y + Inches(0.88), 'Client Panels', BLUE_TINT, SIX_BLUE, w=Inches(1.6), size=11)


# ══════════════════════════════════════════════════════════
# SLIDE 4 — DEMO
# ══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
s4.background.fill.solid()
s4.background.fill.fore_color.rgb = WHITE
topbar(s4, '05 / 06 — LIVE DEMO')

# Left half
tf = tb(s4, Inches(0.8), Inches(1.0), Inches(5.5), Inches(0.3))
para(tf, 'LIVE DEMO', 12, True, SIX_RED)

tf2 = tb(s4, Inches(0.8), Inches(1.5), Inches(5.5), Inches(1.0))
p = tf2.paragraphs[0]
run(p, 'From CRM to advisory ', 32, True, TEXT_DARK)
run(p, 'in 15 seconds', 32, True, SIX_RED, italic=True)

tf3 = tb(s4, Inches(0.8), Inches(2.7), Inches(5.2), Inches(0.7))
para(tf3, 'Four real Swiss private banking personas. Actual CRM logs. Live SIX financial data. One RM. Zero compromise.', 14, False, TEXT_SUB)

# Steps
steps = [
    'Select any client — their DNA surfaces instantly from CRM history',
    'See live conflicts and news matched to their personal values',
    'Hit Generate Advisory — bespoke message in <5s',
    'Ask the AI assistant anything — it knows the client cold',
]
step_y = Inches(3.6)
for i, step in enumerate(steps):
    y = step_y + i * Inches(0.5)
    circ = s4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.3), Inches(0.3))
    circ.fill.solid()
    circ.fill.fore_color.rgb = SIX_RED
    circ.line.fill.background()
    ctf = circ.text_frame
    ctf.margin_top = Pt(0)
    ctf.margin_bottom = Pt(0)
    ctf.margin_left = Pt(0)
    ctf.margin_right = Pt(0)
    cp = ctf.paragraphs[0]
    cp.alignment = CENTER
    run(cp, str(i + 1), 10, True, WHITE)

    tf_s = tb(s4, Inches(1.25), y, Inches(5), Inches(0.4))
    para(tf_s, step, 12, False, TEXT_SUB)

# URL chip
url_box = rect(s4, Inches(0.8), Inches(5.9), Inches(3.2), Inches(0.4), SURFACE, BORDER)
tf_url = tb(s4, Inches(0.95), Inches(5.9), Inches(3), Inches(0.4))
p_url = tf_url.paragraphs[0]
run(p_url, '●  ', 12, False, GREEN)
run(p_url, 'wealthadvisor-ai.fly.dev', 12, False, TEXT_SUB, name='Consolas')

# Right half — mock browser
right_bg = rect(s4, Inches(6.6), Inches(0), Inches(6.733), SH, SURFACE, SURFACE, Pt(0))

browser_x = Inches(7.0)
browser_w = Inches(5.933)
browser_h = Inches(6.2)
browser = rect(s4, browser_x, Inches(0.7), browser_w, browser_h, WHITE, BORDER)

# Browser bar
bar = rect(s4, browser_x, Inches(0.7), browser_w, Inches(0.45), SURFACE, BORDER)
for j, c in enumerate([RGBColor(0xEF,0x44,0x44), RGBColor(0xF5,0x9E,0x0B), RGBColor(0x22,0xC5,0x5E)]):
    dot = s4.shapes.add_shape(MSO_SHAPE.OVAL, browser_x + Inches(0.2 + j * 0.25), Inches(0.83), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = c
    dot.line.fill.background()

url_bar = rect(s4, browser_x + Inches(1.2), Inches(0.82), Inches(3.2), Inches(0.22), WHITE, BORDER)
utf = url_bar.text_frame
utf.margin_top = Pt(0)
utf.margin_bottom = Pt(0)
utf.margin_left = Pt(6)
run(utf.paragraphs[0], 'wealthadvisor-ai.fly.dev', 8, False, TEXT_SUB, name='Consolas')

# Mock sidebar
sb_x = browser_x
sb_w = Inches(1.6)
sidebar = rect(s4, sb_x, Inches(1.15), sb_w, browser_h - Inches(0.45), SURFACE, BORDER)

tf_sb = tb(s4, sb_x + Inches(0.12), Inches(1.2), Inches(1.3), Inches(0.25))
p_sb = tf_sb.paragraphs[0]
run(p_sb, 'SIX ', 9, True, SIX_RED)
run(p_sb, 'AI', 9, False, TEXT_SUB)

clients = [
    ('HS', 'H. Schneider', 'Balanced', SIX_RED, True),
    ('MH', 'M. Huber', 'Defensive', SIX_BLUE, False),
    ('ER', 'E. Räber', 'Defensive', SIX_RED_DK, False),
    ('JA', 'J. Ammann', 'Growth', SIX_BLUE_BR, False),
]
for ci, (init, cname, strat, color, sel) in enumerate(clients):
    cy = Inches(1.6 + ci * 0.5)
    if sel:
        rect(s4, sb_x + Inches(0.05), cy - Inches(0.05), sb_w - Inches(0.1), Inches(0.42), RED_TINT, RED_TINT, Pt(0))

    av = s4.shapes.add_shape(MSO_SHAPE.OVAL, sb_x + Inches(0.12), cy, Inches(0.26), Inches(0.26))
    av.fill.solid()
    av.fill.fore_color.rgb = color
    av.line.fill.background()
    atf = av.text_frame
    atf.margin_top = Pt(0)
    atf.margin_bottom = Pt(0)
    ap = atf.paragraphs[0]
    ap.alignment = CENTER
    run(ap, init, 7, True, WHITE)

    tf_cn = tb(s4, sb_x + Inches(0.45), cy - Inches(0.03), Inches(1.1), Inches(0.35))
    para(tf_cn, cname, 8, True, TEXT_DARK)
    para(tf_cn, strat, 7, False, TEXT_SUB)

# Mock main content
main_x = sb_x + sb_w + Inches(0.1)
main_w = browser_w - sb_w - Inches(0.1)

# Advisory
adv_box = rect(s4, main_x + Inches(0.15), Inches(1.25), main_w - Inches(0.3), Inches(1.2), RED_TINT, RGBColor(0xF0,0xD0,0xD5))
tf_adv = tb(s4, main_x + Inches(0.3), Inches(1.3), main_w - Inches(0.6), Inches(1.1))
p_adv = tf_adv.paragraphs[0]
run(p_adv, 'Advisory Draft — ', 8, True, SIX_RED)
run(p_adv, "Dear Hubertus, the announced shutdown of AstraZeneca's chronic-illness research division directly conflicts with your foundation's mission. We recommend a swap into Roche Holding.", 8, False, TEXT_SUB)

# DNA card
dna_box = rect(s4, main_x + Inches(0.15), Inches(2.55), main_w - Inches(0.3), Inches(0.85), SURFACE, BORDER)
tf_dna = tb(s4, main_x + Inches(0.3), Inches(2.58), Inches(2), Inches(0.25))
para(tf_dna, 'CLIENT DNA', 8, True, SIX_RED)

dna_pills_data = [('health legacy', BLUE_TINT, SIX_BLUE), ('ESG alignment', BLUE_TINT, SIX_BLUE),
                   ('pharma risk', RED_TINT, SIX_RED), ('family foundation', PURPLE_TINT, PURPLE)]
dpx = main_x + Inches(0.3)
for pt, bg, tc in dna_pills_data:
    pw = Inches(max(0.9, len(pt) * 0.1))
    pill(s4, dpx, Inches(2.9), pt, bg, tc, w=pw, h=Inches(0.28), size=8)
    dpx += pw + Inches(0.1)

# Alert
alert_box = rect(s4, main_x + Inches(0.15), Inches(3.5), main_w - Inches(0.3), Inches(0.5), RED_TINT, RGBColor(0xF0,0xD0,0xD5))
tf_al = tb(s4, main_x + Inches(0.3), Inches(3.52), main_w - Inches(0.6), Inches(0.45))
para(tf_al, '⚠ Major Pharma Co. shuts chronic-illness research — direct conflict with client DNA', 8, False, SIX_RED_DK)


# ══════════════════════════════════════════════════════════
# SLIDE 5 — Q&A / TEAM
# ══════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
s5.background.fill.solid()
s5.background.fill.fore_color.rgb = WHITE
topbar(s5, '06 / 06 — Q&A')

# Eyebrow
tf = tb(s5, LM, Inches(0.85), CONTENT_W, Inches(0.3))
para(tf, 'THE TEAM', 12, True, SIX_RED, CENTER)

# Title
tf2 = tb(s5, Inches(2), Inches(1.2), Inches(9.333), Inches(0.7))
p = tf2.paragraphs[0]
p.alignment = CENTER
run(p, 'Questions & ', 36, True, TEXT_DARK)
run(p, 'Answers', 36, True, SIX_RED, italic=True)

# Subtitle
tf3 = tb(s5, Inches(2), Inches(1.9), Inches(9.333), Inches(0.4))
para(tf3, 'Built in ~24 hours at SwissHacks 2026', 15, False, TEXT_SUB, CENTER)

# Team cards
team = [
    ('boris.JPG', 'Boris Goranov', 'FRONT-END', SIX_BLUE, BLUE_TINT),
    ('hristo.jpeg', 'Hristo Stefanov', 'BACK-END', GREEN, GREEN_TINT),
    ('nikola.jpeg', 'Nikola Staykov', 'BUSINESS', AMBER, AMBER_TINT),
]

card_w = Inches(3.2)
card_h = Inches(3.4)
card_gap = Inches(0.5)
total_cards_w = 3 * card_w.inches + 2 * card_gap.inches
card_start_x = Inches((13.333 - total_cards_w) / 2)
card_y = Inches(2.5)

for i, (photo, name, role, role_color, role_bg) in enumerate(team):
    x = Inches(card_start_x.inches + i * (card_w.inches + card_gap.inches))
    box = rect(s5, x, card_y, card_w, card_h)

    # Photo
    photo_path = os.path.join(DIR, photo)
    if os.path.exists(photo_path):
        photo_size = Inches(1.8)
        photo_x = Inches(x.inches + (card_w.inches - photo_size.inches) / 2)
        s5.shapes.add_picture(photo_path, photo_x, card_y + Inches(0.3), photo_size, photo_size)

    # Name
    tf_name = tb(s5, x, card_y + Inches(2.2), card_w, Inches(0.4))
    para(tf_name, name, 16, True, TEXT_DARK, CENTER)

    # Role pill
    role_w = Inches(max(1.2, len(role) * 0.14))
    role_x = Inches(x.inches + (card_w.inches - role_w.inches) / 2)
    pill(s5, role_x, card_y + Inches(2.7), role, role_bg, role_color, w=role_w, size=11)

# Closing statement
tf_close = tb(s5, Inches(2), Inches(6.15), Inches(9.333), Inches(0.5))
para(tf_close, 'Every client deserves a personal banker. We built one.', 20, True, TEXT_DARK, CENTER)

# Questions
tf_q = tb(s5, Inches(4), Inches(6.7), Inches(5.333), Inches(0.4))
para(tf_q, '\U0001F4AC  Questions?', 16, True, TEXT_SUB, CENTER)


# ── Save ────────────────────────────────────────────────
out_path = os.path.join(DIR, 'SIX_AI_Presentation.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
